"""Extract plain text from Office documents with sha256-keyed caching.

**Scope (updated 2026-04-18):** this script handles only file types where a
pure-Python library already gives reliable text: .docx / .pptx / .xlsx.

Everything else is out of scope by design:

- .md / .txt / .csv / .json → the Skill reads these directly via the platform's
  Read tool (no extraction needed).
- .pdf / .jpg / .png / .jpeg / .webp → the Skill reads these via Claude's
  native multimodal Read (the model ingests them as images/PDF pages).
  Trying to OCR or extract here would be strictly worse than letting the
  model see the original.

Use `is_supported(path)` to check before calling `extract_text(path, cache_dir)`.
Unsupported types raise UnsupportedExtension — the Skill layer routes around
them instead of catching blindly.
"""

import hashlib
from pathlib import Path


SUPPORTED_EXTS: set[str] = {".docx", ".pptx", ".xlsx"}


class UnsupportedExtension(ValueError):
    """Raised when extract_text is called on a file it cannot handle.

    The Skill should route PDFs/images/markdown to the appropriate reader
    rather than catching this — it's a routing bug, not a runtime condition.
    """


def is_supported(path: Path) -> bool:
    """True if extract_text can handle this file. Use this for routing."""
    return path.suffix.lower() in SUPPORTED_EXTS


def _sha256(p: Path) -> str:
    h = hashlib.sha256()
    with p.open("rb") as f:
        for chunk in iter(lambda: f.read(1 << 16), b""):
            h.update(chunk)
    return h.hexdigest()


def cache_path(src: Path, cache_dir: Path) -> Path:
    return cache_dir / f"{_sha256(src)}.txt"


def get_cached(src: Path, cache_dir: Path) -> str | None:
    p = cache_path(src, cache_dir)
    return p.read_text(encoding="utf-8") if p.exists() else None


def extract_text(src: Path, cache_dir: Path) -> str:
    """Extract text from a supported Office file. Uses sha256 cache.

    Raises UnsupportedExtension for anything other than .docx/.pptx/.xlsx.
    """
    cache_dir.mkdir(parents=True, exist_ok=True)
    hit = get_cached(src, cache_dir)
    if hit is not None:
        return hit

    ext = src.suffix.lower()
    if ext == ".docx":
        text = _extract_docx(src)
    elif ext == ".pptx":
        text = _extract_pptx(src)
    elif ext == ".xlsx":
        text = _extract_xlsx(src)
    else:
        raise UnsupportedExtension(
            f"{src.name}: extract_text only supports .docx/.pptx/.xlsx. "
            f"PDFs and images are read natively by the Skill; .md/.txt are read directly."
        )

    cache_path(src, cache_dir).write_text(text, encoding="utf-8")
    return text


def _extract_docx(src: Path) -> str:
    from docx import Document
    doc = Document(str(src))
    parts: list[str] = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_pptx(src: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(src))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"--- Slide {i} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    t = "".join(r.text for r in p.runs)
                    if t.strip():
                        parts.append(t)
    return "\n".join(parts)


def _extract_xlsx(src: Path) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(str(src), data_only=True)
    parts: list[str] = []
    for ws in wb.worksheets:
        parts.append(f"--- Sheet: {ws.title} ---")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)
